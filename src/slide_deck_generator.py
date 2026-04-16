from __future__ import annotations

import pathlib
import re
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .topic_research import TopicResearchBrief

# ── Colour palette ───────────────────────────────────────────────────────────
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1C, 0x1C, 0x2E)  # dark navy — section divider backgrounds
C_SURFACE = RGBColor(0xF5, 0xF4, 0xF1)  # warm off-white surface
C_TEXT    = RGBColor(0x1F, 0x1F, 0x1F)  # near-black body text
C_SUBTLE  = RGBColor(0x78, 0x77, 0x74)  # muted grey labels
C_ACCENT  = RGBColor(0x23, 0x83, 0xE2)  # primary blue
C_TEAL    = RGBColor(0x0D, 0x94, 0x88)  # teal — guiding-question slides
C_DIVIDER = RGBColor(0xE8, 0xE7, 0xE4)  # hairline dividers

# ── Typography  (Futura; falls back to Century Gothic / sans-serif) ──────────
FONT_HEADING = "Futura"
FONT_BODY    = "Futura"

# ── Layout constants (13.33 × 7.5 in, standard widescreen) ──────────────────
W         = Inches(13.33)
H         = Inches(7.5)
MARGIN_L  = Inches(0.9)
MARGIN_R  = Inches(0.9)
MARGIN_T  = Inches(0.6)
CONTENT_W = W - MARGIN_L - MARGIN_R   # ≈ 11.53 in


class SlideDeckGenerator:
    def __init__(self, output_path: pathlib.Path):
        self.output_path = output_path
        self.presentation = Presentation()
        self.presentation.slide_width  = W
        self.presentation.slide_height = H
        self._blank = self._find_blank_layout()
        self._slide_num = 0

    # ── Internal helpers ─────────────────────────────────────────────────────

    def _find_blank_layout(self):
        for layout in self.presentation.slide_layouts:
            if layout.name.lower() == "blank":
                return layout
        return self.presentation.slide_layouts[6]

    def _new_slide(self):
        slide = self.presentation.slides.add_slide(self._blank)
        # Remove any placeholder shapes that the blank layout might carry
        for ph in list(slide.placeholders):
            ph._element.getparent().remove(ph._element)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_WHITE
        self._slide_num += 1
        return slide

    def _set_bg(self, slide, color: RGBColor):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _txb(self, slide, left, top, width, height):
        """Add a textbox and return its text_frame with word-wrap on."""
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        return tf

    def _rect(self, slide, left, top, width, height, color: RGBColor):
        """Add a solid filled, borderless rectangle."""
        shp = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def _set_para(self, p, text: str, font: str, size: float, color: RGBColor,
                  bold: bool = False, italic: bool = False,
                  align: PP_ALIGN = PP_ALIGN.LEFT,
                  space_before: float = 0, space_after: float = 0):
        p.text = text
        p.alignment = align
        if space_before:
            p.space_before = Pt(space_before)
        if space_after:
            p.space_after = Pt(space_after)
        for run in p.runs:
            run.font.name   = font
            run.font.size   = Pt(size)
            run.font.color.rgb = color
            run.font.bold   = bold
            run.font.italic = italic

    def _add_para(self, tf, text: str, font: str, size: float, color: RGBColor,
                  bold: bool = False, space_before: float = 0,
                  align: PP_ALIGN = PP_ALIGN.LEFT):
        p = tf.add_paragraph()
        self._set_para(p, text, font, size, color, bold=bold,
                       space_before=space_before, align=align)
        return p

    def _attach_notes(self, slide, notes_text: str):
        slide.notes_slide.notes_text_frame.text = notes_text

    def _trim(self, text: str) -> str:
        return text.strip()

    def _prose_to_bullets(self, text: str, max_bullets: int = 6) -> List[str]:
        """Split prose into individual sentences, one per bullet point."""
        bullets = []
        for para in text.split("\n\n"):
            para = para.strip()
            if not para:
                continue
            sentences = re.split(r'(?<=[.!?])\s+', para)
            for s in sentences:
                s = s.strip()
                if s:
                    bullets.append(self._trim(s))
                if len(bullets) >= max_bullets:
                    return bullets
        return bullets

    # ── Public slide builders ────────────────────────────────────────────────

    def add_title_slide(self, course_title: str, author: str = "Course Agent"):
        slide = self._new_slide()

        # Left accent spine — full-height bar
        SPINE_W = Inches(0.42)
        self._rect(slide, Inches(0), Inches(0), SPINE_W, H, C_ACCENT)

        # Decorative surface block — bottom-right corner
        self._rect(slide, W - Inches(4.8), H - Inches(2.0), Inches(4.8), Inches(2.0), C_SURFACE)

        # Small accent dot — sits between spine and title area, vertically centred
        DOT = Inches(0.10)
        self._rect(slide, SPINE_W + Inches(0.18), Inches(2.05), DOT, DOT, C_ACCENT)

        # Title
        TITLE_L = SPINE_W + Inches(0.38)
        TITLE_W = W - TITLE_L - Inches(0.5)
        tf = self._txb(slide, TITLE_L, Inches(2.2), TITLE_W, Inches(2.0))
        self._set_para(tf.paragraphs[0], course_title,
                       FONT_HEADING, 38, C_TEXT, bold=False)

        # Author
        tf2 = self._txb(slide, TITLE_L, Inches(4.45), TITLE_W, Inches(0.7))
        self._set_para(tf2.paragraphs[0], f"Prepared by {author}",
                       FONT_BODY, 15, C_SUBTLE)

        # Bottom bar (starts after spine)
        self._rect(slide, SPINE_W, H - Inches(0.48), W - SPINE_W, Inches(0.48), C_SURFACE)
        tf3 = self._txb(slide, TITLE_L, H - Inches(0.44), TITLE_W, Inches(0.38))
        self._set_para(tf3.paragraphs[0], "Research Studio", FONT_BODY, 10, C_SUBTLE)

        notes = self._build_speaker_notes(
            title=course_title,
            bullets=[f"Instructor: {author}",
                     "Introduce learning outcomes and expected participation.",
                     "Set context for the module progression and timing."],
            slide_type="title",
        )
        self._attach_notes(slide, notes)

    def add_section_slide(self, title: str, bullets: List[str], notes_text: Optional[str] = None):
        """Add a content slide. Section-header titles get a distinctive look."""
        clean = [b.strip() for b in bullets if b and b.strip()]
        slide_type = self._infer_slide_type(title)

        if slide_type in ("section-divider",) or re.match(r"^Part \d+", title):
            self._build_section_divider(title, clean)
        else:
            self._build_content_slide(title, clean, slide_type, notes_text=notes_text)

    def _build_section_divider(self, title: str, bullets: List[str]):
        slide = self._new_slide()
        self._set_bg(slide, C_DARK)

        # Accent right-edge strip
        self._rect(slide, W - Inches(0.42), Inches(0), Inches(0.42), H, C_ACCENT)

        # Subtle horizontal rule above title
        self._rect(slide, MARGIN_L, Inches(2.85), Inches(1.6), Pt(2),
                   RGBColor(0x4A, 0x4A, 0x6A))

        # Section title — large, white
        tf = self._txb(slide, MARGIN_L, Inches(3.0), CONTENT_W - Inches(0.6), Inches(2.0))
        self._set_para(tf.paragraphs[0], title, FONT_HEADING, 34, C_WHITE)

        # Optional first bullet as a light descriptor
        if bullets:
            tf2 = self._txb(slide, MARGIN_L, Inches(4.75), CONTENT_W - Inches(0.6), Inches(0.8))
            self._set_para(tf2.paragraphs[0], bullets[0], FONT_BODY, 15,
                           RGBColor(0x9A, 0x99, 0x96))

        notes = self._build_speaker_notes(title, bullets, "generic")
        self._attach_notes(slide, notes)

    def _build_content_slide(self, title: str, bullets: List[str], slide_type: str,
                             notes_text: Optional[str] = None,
                             accent_color: Optional[RGBColor] = None):
        slide = self._new_slide()
        color = accent_color or C_ACCENT

        # Thin top accent strip
        self._rect(slide, Inches(0), Inches(0), W, Inches(0.07), color)

        # Left-margin tick mark aligned with title
        self._rect(slide, Inches(0.38), MARGIN_T + Inches(0.04), Inches(0.07), Inches(0.72), color)

        # Title
        tf = self._txb(slide, MARGIN_L, MARGIN_T, CONTENT_W, Inches(0.85))
        self._set_para(tf.paragraphs[0], title, FONT_HEADING, 24, C_TEXT)

        # Hairline divider under title
        self._rect(slide, MARGIN_L, Inches(1.45), CONTENT_W, Pt(1), C_DIVIDER)

        # Body
        if bullets:
            BODY_T = Inches(1.62)
            BODY_H = H - BODY_T - Inches(0.5)
            tf2 = self._txb(slide, MARGIN_L, BODY_T, CONTENT_W, BODY_H)
            self._set_para(tf2.paragraphs[0], f"• {bullets[0]}", FONT_BODY, 14, C_TEXT)
            for bullet in bullets[1:]:
                self._add_para(tf2, f"• {bullet}", FONT_BODY, 14, C_TEXT,
                               space_before=6)

        # Bottom bar with slide number
        self._rect(slide, Inches(0), H - Inches(0.38), W, Inches(0.38), C_SURFACE)
        tf3 = self._txb(slide, MARGIN_L, H - Inches(0.35), CONTENT_W - Inches(1.0), Inches(0.3))
        self._set_para(tf3.paragraphs[0], "Research Studio", FONT_BODY, 9, C_SUBTLE)
        tf4 = self._txb(slide, W - Inches(1.3), H - Inches(0.35), Inches(1.0), Inches(0.3))
        self._set_para(tf4.paragraphs[0], str(self._slide_num), FONT_BODY, 9, C_SUBTLE,
                       align=PP_ALIGN.RIGHT)

        notes = notes_text if notes_text else self._build_speaker_notes(title, bullets, slide_type)
        self._attach_notes(slide, notes)

    # ── Slide type inference ─────────────────────────────────────────────────

    def _infer_slide_type(self, title: str) -> str:
        lowered = title.lower()
        if re.match(r"^part \d+", lowered):
            return "section-divider"
        if "overview" in lowered or "roadmap" in lowered:
            return "overview"
        if "objective" in lowered or "prerequisite" in lowered:
            return "course-setup"
        if "foundation" in lowered or "core concept" in lowered or "historical" in lowered:
            return "foundations"
        if "method" in lowered or "model" in lowered or "evaluation" in lowered:
            return "methods"
        if "application" in lowered or "domain" in lowered:
            return "applications"
        if "case study" in lowered:
            return "case-study"
        if "deployment" in lowered or "monitor" in lowered:
            return "deployment"
        if "ethical" in lowered or "accountability" in lowered:
            return "ethics"
        if "pitfall" in lowered or "misconception" in lowered or "failure" in lowered:
            return "pitfalls"
        if "lab" in lowered or "exercise" in lowered:
            return "lab"
        if "discussion" in lowered or "capstone" in lowered:
            return "discussion"
        if "reference" in lowered:
            return "references"
        return "generic"

    # ── Speaker notes ────────────────────────────────────────────────────────

    def _build_speaker_notes(self, title: str, bullets: List[str], slide_type: str = "generic") -> str:
        lines = [title, ""]
        for i, b in enumerate(bullets, 1):
            lines.append(f"{i}. {b}")
        return "\n".join(lines)

    # ── Research-specific slides ─────────────────────────────────────────────

    def add_topic_research_slides(self, idx: int, brief: TopicResearchBrief):
        deep_dive_bullets = [
            f"Summary: {brief.summary[:220]}",
            f"Source score: {brief.source_score:.2f} | Citation confidence: {brief.citation_confidence}",
            f"Talking point 1: {brief.talking_points[0]}",
            f"Talking point 2: {brief.talking_points[1]}",
            f"Talking point 3: {brief.talking_points[2]}",
        ]
        self.add_section_slide(f"Module {idx} Deep Dive: {brief.topic}", deep_dive_bullets)

        discussion_bullets = [
            f"Prompt 1: {brief.discussion_questions[0]}",
            f"Prompt 2: {brief.discussion_questions[1]}",
            f"Prompt 3: {brief.discussion_questions[2]}",
            "Ask students to support their claims with evidence from the source readings.",
        ]
        self.add_section_slide(f"Module {idx} Discussion Lab", discussion_bullets)

    def add_guiding_question_slide(self, question: str, response: str, bullets: Optional[List[str]] = None):
        """One slide per guiding question — concise bullets on slide, full prose in notes."""
        slide = self._new_slide()

        # Thin top accent strip (teal to distinguish from regular slides)
        self._rect(slide, Inches(0), Inches(0), W, Inches(0.07), C_TEAL)

        # Left-margin tick mark
        self._rect(slide, Inches(0.38), MARGIN_T + Inches(0.04), Inches(0.07), Inches(1.0), C_TEAL)

        # Question as title
        tf = self._txb(slide, MARGIN_L, MARGIN_T, CONTENT_W, Inches(1.1))
        self._set_para(tf.paragraphs[0], question, FONT_HEADING, 20, C_TEXT)

        # Teal accent divider under question
        self._rect(slide, MARGIN_L, Inches(1.62), CONTENT_W, Pt(2), C_TEAL)

        # Use LLM-generated concise bullets if available, else extract from prose
        bullets = bullets or self._prose_to_bullets(response, max_bullets=6)
        if not bullets:
            bullets = [self._trim(response)]

        BODY_T = Inches(1.78)
        BODY_H = H - BODY_T - Inches(0.5)
        tf2 = self._txb(slide, MARGIN_L, BODY_T, CONTENT_W, BODY_H)
        self._set_para(tf2.paragraphs[0], f"• {bullets[0]}", FONT_BODY, 14, C_TEXT)
        for b in bullets[1:]:
            self._add_para(tf2, f"• {b}", FONT_BODY, 14, C_TEXT, space_before=6)

        # Bottom bar with slide number
        self._rect(slide, Inches(0), H - Inches(0.38), W, Inches(0.38), C_SURFACE)
        tf3 = self._txb(slide, MARGIN_L, H - Inches(0.35), CONTENT_W - Inches(1.0), Inches(0.3))
        self._set_para(tf3.paragraphs[0], "Research Studio", FONT_BODY, 9, C_SUBTLE)
        tf4 = self._txb(slide, W - Inches(1.3), H - Inches(0.35), Inches(1.0), Inches(0.3))
        self._set_para(tf4.paragraphs[0], str(self._slide_num), FONT_BODY, 9, C_SUBTLE,
                       align=PP_ALIGN.RIGHT)

        # Full prose goes into speaker notes for presenter reference
        self._attach_notes(slide, f"Question: {question}\n\nFull answer:\n{response}")

    def add_references_slide(self, research_brief: TopicResearchBrief):
        source_lines: List[str] = []
        if research_brief.citations:
            for citation in research_brief.citations[:6]:
                source_lines.append(
                    f"{citation.source} [{citation.citation_confidence} {citation.weighted_score:.2f}]: {citation.url}"
                )
        else:
            for source in research_brief.sources:
                source_lines.append(f"Reference: {source}")

        self.add_section_slide("Research References", source_lines or ["No external references collected."])

    def save(self):
        self.presentation.save(self.output_path)

    # ── Full deck builders ───────────────────────────────────────────────────

    def build_research_deck(self, topic: str, research_brief: TopicResearchBrief, course_title: str, author: str):
        """Generate a lean, LLM-driven research deck — every slide comes from actual content."""
        tp             = research_brief.talking_points or []
        objectives     = research_brief.learning_objectives or []
        prerequisites  = research_brief.prerequisite_topics or []
        case_studies   = research_brief.case_studies or []
        labs           = research_brief.lab_exercises or []
        misconceptions = research_brief.misconceptions or []
        discussion     = research_brief.discussion_questions or []
        modules        = research_brief.curriculum_modules or {}

        # ── Title ────────────────────────────────────────────────────────────
        self.add_title_slide(topic, author)

        # ── Overview: bullets on slide, full summary prose in notes ──────────
        summary_bullets = self._prose_to_bullets(research_brief.summary, max_bullets=4)
        if summary_bullets:
            self.add_section_slide("Overview", summary_bullets,
                                   notes_text=research_brief.summary)

        # ── Research Findings: one slide per guiding question ─────────────────
        gqr = getattr(research_brief, "guiding_question_responses", [])
        if gqr:
            self._build_section_divider(
                "Research Findings",
                [f"{len(gqr)} guiding question{'s' if len(gqr) != 1 else ''} answered"],
            )
            for item in gqr:
                self.add_guiding_question_slide(
                    question=item.get("question", ""),
                    response=item.get("response", ""),
                    bullets=item.get("bullets") or None,
                )

        # ── Key insights from talking points ──────────────────────────────────
        if tp:
            self.add_section_slide("Key Insights", [self._trim(t) for t in tp[:5]])

        # ── Curriculum modules — one slide per section ────────────────────────
        MODULE_LABELS = {
            "foundations":  "Foundations",
            "methods":      "Methods & Modeling",
            "applications": "Applications",
            "evaluation":   "Evaluation",
            "advanced":     "Advanced Topics",
        }
        for key, label in MODULE_LABELS.items():
            items = [i for i in (modules.get(key) or []) if i]
            if items:
                self.add_section_slide(label, [self._trim(i) for i in items[:4]])

        # ── Learning objectives ───────────────────────────────────────────────
        if objectives:
            self.add_section_slide("Learning Objectives", [self._trim(o) for o in objectives[:5]])

        # ── Prerequisites ─────────────────────────────────────────────────────
        if prerequisites:
            self.add_section_slide("Prerequisites", [self._trim(p) for p in prerequisites[:5]])

        # ── Case studies ──────────────────────────────────────────────────────
        if case_studies:
            self.add_section_slide("Case Studies", [self._trim(c) for c in case_studies[:4]])

        # ── Common misconceptions ─────────────────────────────────────────────
        if misconceptions:
            self.add_section_slide("Common Misconceptions", [self._trim(m) for m in misconceptions[:4]])

        # ── Hands-on exercises ────────────────────────────────────────────────
        if labs:
            self.add_section_slide("Hands-on Exercises", [self._trim(l) for l in labs[:4]])

        # ── Discussion questions ──────────────────────────────────────────────
        if discussion:
            self.add_section_slide("Discussion Questions", [self._trim(d) for d in discussion[:4]])

        # ── References ───────────────────────────────────────────────────────
        self.add_references_slide(research_brief)

    def build_deck(self, course_title: str, topics: List[str], source_summaries: List[str],
                   research_briefs: Optional[List[TopicResearchBrief]] = None):
        """Legacy method for backward compatibility."""
        self.add_title_slide(course_title)
        self.add_section_slide("Course Overview", [
            "3-hour session split into modules",
            f"Core topics: {', '.join(topics)}",
            "Interactive lecture with researched talking points and discussion prompts",
        ])

        brief_by_topic = {brief.topic: brief for brief in (research_briefs or [])}

        for idx, topic in enumerate(topics, start=1):
            bullets = [
                f"What is {topic}?",
                f"Key concepts for {topic} from course materials",
                "Applied examples and practice problems",
                "Transition into evidence-based discussion prompts",
            ]
            self.add_section_slide(f"Module {idx}: {topic}", bullets)

            brief = brief_by_topic.get(topic)
            if brief:
                self.add_topic_research_slides(idx, brief)

        self.add_section_slide("Source Content Summary", source_summaries[:6] or ["No source content identified."])

        if research_briefs:
            self.add_references_slide(research_briefs[0])
